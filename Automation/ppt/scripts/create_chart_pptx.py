# -*- coding: utf-8 -*-
"""
Excel → PPTX 편집 가능 차트 (v4 — 흰 배경 + 안정)
- 배경: 흰색 (워드 프로세서 삽입용)
- 라벨: 검정/진회색
- python-pptx 표준 API만 사용
"""

import sys, os, re
from datetime import datetime
from dataclasses import dataclass
from typing import Tuple
sys.stdout.reconfigure(encoding='utf-8', errors='replace')

import openpyxl
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData


# ═══════════════════════════════════════════════════════
# 테마 (흰 배경용)
# ═══════════════════════════════════════════════════════

@dataclass
class Theme:
    # 차트 색상
    accent: Tuple = (0x3B, 0x82, 0xF6)     # 파란색
    accent2: Tuple = (0x10, 0xB9, 0x81)     # 녹색
    accent3: Tuple = (0xF5, 0x9E, 0x0B)     # 주황
    accent4: Tuple = (0x8B, 0x5C, 0xF6)     # 보라
    accent5: Tuple = (0xEF, 0x44, 0x44)     # 빨강
    # 텍스트 색상
    black: Tuple = (0x22, 0x22, 0x22)       # 제목/라벨
    dark: Tuple = (0x33, 0x33, 0x33)        # 본문
    gray: Tuple = (0x66, 0x66, 0x66)        # 보조
    light_gray: Tuple = (0xCC, 0xCC, 0xCC)  # 축선
    white: Tuple = (0xFF, 0xFF, 0xFF)
    font: str = "맑은 고딕"
    @property
    def colors(self):
        return [self.accent, self.accent2, self.accent3, self.accent4, self.accent5,
                (0x06,0xB6,0xD4),(0xF9,0x73,0x16),(0xEC,0x48,0x99),(0x14,0xB8,0xA6),(0xA8,0x55,0xF7)]
    def rgb(self, t): return RGBColor(*t)

T = Theme()


# ═══════════════════════════════════════════════════════
# 표 설정
# ═══════════════════════════════════════════════════════

@dataclass
class Cfg:
    idx: int; title: str; total_row: int; name_row: int
    data_col: int = 4; n_col: int = 3
    chart: str = "hbar"; end_row: int = 0

TABLES = [
    Cfg(1,  "응답자 일반현황",                          3,   4, data_col=3, chart="pie"),
    Cfg(2,  "일자리 정보 탐색 활용 경로 (중복응답)",    21,  19),
    Cfg(3,  "취업 준비 중 느끼는 어려움 (중복응답)",    39,  37),
    Cfg(4,  "대학에서 제공 받았으면 하는 지원 (중복응답)", 57, 55),
    Cfg(5,  "필요 취업 정보 (1순위)",                   75,  73),
    Cfg(6,  "필요 취업 정보 (1순위+2순위)",             93,  91),
    Cfg(7,  "부족 역량 (1순위)",                        111, 109),
    Cfg(8,  "부족 역량 (1순위+2순위)",                  129, 127),
    Cfg(9,  "취업 시 고려 요소 (1순위)",                147, 145),
    Cfg(10, "취업 시 고려 요소 (1순위+2순위)",          165, 163),
    Cfg(11, "대학교별 취업 희망 업종",                  183, 181, chart="stacked", end_row=226),
    Cfg(12, "대학교별 취업 희망 직군",                  232, 230, chart="stacked", end_row=267),
    Cfg(13, "울산 지역 기업 인식 수준 (평균 점수)",     289, 288, chart="radar"),
    Cfg(14, "졸업 후 희망 취업 지역",                   307, 305, chart="pie"),
    Cfg(15, "울산 정주환경이 타지역 취업에 미친 영향",  325, 323, chart="donut"),
    Cfg(16, "타지역 취업에 영향을 미친 정주환경 요인 (중복응답)", 343, 341),
    Cfg(17, "울산 지역 청년 정착 정책 (1순위)",         361, 359),
    Cfg(18, "울산 지역 청년 정착 정책 (1순위+2순위)",   379, 377),
]


# ═══════════════════════════════════════════════════════
# 데이터 추출
# ═══════════════════════════════════════════════════════

def get_bar_data(ws, c):
    labels, pcts, freqs = [], [], []
    n = int(ws.cell(row=c.total_row, column=c.n_col).value or 0)
    for col in range(c.data_col, ws.max_column + 1, 2):
        nm = ws.cell(row=c.name_row, column=col).value
        if nm is None or str(nm).strip() == '': continue
        s = str(nm).strip()
        if s in ('빈도', '비율', '구 분'): continue
        if len(s) > 22: s = s[:20] + "…"
        labels.append(s)
        fv = ws.cell(row=c.total_row, column=col).value
        freqs.append(float(fv) if fv else 0)
        pv = ws.cell(row=c.total_row, column=col + 1).value
        pcts.append(round(float(pv) * 100, 1) if pv else 0)
    return labels, pcts, freqs, n

def get_pie1(ws):
    labels, vals = [], []
    for r in range(4, 8):
        nm = ws.cell(row=r, column=2).value
        v  = ws.cell(row=r, column=3).value
        if nm and v:
            labels.append(str(nm).strip())
            vals.append(float(v))
    total = sum(vals)
    pcts = [round(v / total * 100, 1) for v in vals]
    return labels, pcts, vals, int(total)

def get_radar(ws, c):
    labels, vals = [], []
    for col in range(3, ws.max_column + 1):
        nm = ws.cell(row=c.name_row, column=col).value
        if nm is None: continue
        v = ws.cell(row=c.total_row, column=col).value
        if v is not None:
            s = str(nm).strip()
            if len(s) > 14: s = s[:12] + "…"
            labels.append(s)
            vals.append(round(float(v), 2))
    return labels, vals

def get_stacked(ws, c, top_n=10):
    univs = []
    for col in range(3, 7):
        v = ws.cell(row=c.name_row, column=col).value
        if v: univs.append(str(v).strip())
    if not univs:
        univs = ["울산대학교", "울산과학대학교", "춘해보건대학교", "울산과학기술원"]
    items, data = [], {u: [] for u in univs}
    for row in range(c.total_row + 1, c.end_row + 1 if c.end_row else c.total_row + 80):
        a = ws.cell(row=row, column=1).value
        if a is None or str(a).strip() == '': break
        a_s = str(a).strip()
        if a_s in ('전 체', '전체'): continue
        m = re.match(r'^\d+\s+(.+)$', a_s)
        name = m.group(1) if m else a_s
        if len(name) > 20: name = name[:18] + "…"
        bv = ws.cell(row=row, column=2).value
        if bv is None or float(bv) == 0: continue
        items.append(name)
        for ci, u in enumerate(univs):
            v = ws.cell(row=row, column=3 + ci).value
            data[u].append(float(v) if v else 0)
    totals = [sum(data[u][i] for u in univs) for i in range(len(items))]
    idx_s = sorted(range(len(items)), key=lambda i: totals[i], reverse=True)[:top_n]
    return [items[i] for i in idx_s], {u: [data[u][i] for i in idx_s] for u in univs}, univs


# ═══════════════════════════════════════════════════════
# 축 투명 처리
# ═══════════════════════════════════════════════════════

def _hide_value_axis(chart):
    va = chart.value_axis
    va.has_major_gridlines = False
    va.has_minor_gridlines = False
    va.has_title = False
    va.tick_labels.font.size = Pt(2)
    va.tick_labels.font.color.rgb = T.rgb(T.white)  # 흰 배경과 동일
    va.format.line.color.rgb = T.rgb(T.white)
    va.format.line.width = Pt(0)


# ═══════════════════════════════════════════════════════
# 차트 빌더 (흰 배경)
# ═══════════════════════════════════════════════════════

class Builder:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

    def _slide(self):
        """흰 배경 슬라이드 — 장식 없음"""
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        # 배경 흰색 (기본값이지만 명시)
        bg = s.background.fill
        bg.solid()
        bg.fore_color.rgb = T.rgb(T.white)
        return s

    def _title(self, slide, text, sub=""):
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(10), Inches(0.5))
        p = tb.text_frame.paragraphs[0]
        r = p.add_run(); r.text = text
        r.font.size = Pt(22); r.font.color.rgb = T.rgb(T.black)
        r.font.bold = True; r.font.name = T.font
        if sub:
            r2 = p.add_run(); r2.text = f"  {sub}"
            r2.font.size = Pt(12); r2.font.color.rgb = T.rgb(T.gray); r2.font.name = T.font

    def _unit(self, slide, text):
        tb = slide.shapes.add_textbox(Inches(10.5), Inches(0.35), Inches(2.5), Inches(0.35))
        p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
        r = p.add_run(); r.text = f"(단위 : {text})"
        r.font.size = Pt(10); r.font.color.rgb = T.rgb(T.gray); r.font.name = T.font

    def _pgnum(self, slide, n, total):
        tb = slide.shapes.add_textbox(Inches(12), Inches(7.05), Inches(1), Inches(0.3))
        p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
        r = p.add_run(); r.text = f"{n}/{total}"
        r.font.size = Pt(10); r.font.color.rgb = T.rgb(T.gray); r.font.name = T.font

    # ── 수평 막대 ──
    def add_hbar(self, labels, pcts, freqs, n, title, pg, tot):
        slide = self._slide()
        self._title(slide, title, f"(n={n:,})")
        self._unit(slide, "%, 명")

        pairs = sorted(zip(labels, pcts, freqs), key=lambda x: x[1])
        s_l, s_p, s_f = [p[0] for p in pairs], [p[1] for p in pairs], [p[2] for p in pairs]

        cd = CategoryChartData()
        cd.categories = s_l
        cd.add_series(' ', s_p)

        cf = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED,
            Inches(0.5), Inches(1.0), Inches(12.3), Inches(5.9), cd
        )
        chart = cf.chart
        chart.has_legend = False

        ser = chart.series[0]
        ser.format.fill.solid()
        ser.format.fill.fore_color.rgb = T.rgb(T.accent)

        # 데이터 라벨
        plot = chart.plots[0]
        plot.has_data_labels = True
        dl = plot.data_labels
        dl.font.size = Pt(10)
        dl.font.color.rgb = T.rgb(T.dark)
        dl.font.name = T.font
        dl.number_format = '0.0"%"'
        dl.label_position = XL_LABEL_POSITION.OUTSIDE_END

        # 커스텀 라벨 (비율% + 빈도)
        for i in range(len(s_l)):
            pt_dl = ser.points[i].data_label
            tf = pt_dl.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = f"{s_p[i]:.1f}% ({int(s_f[i])})"
            run.font.size = Pt(10)
            run.font.color.rgb = T.rgb(T.dark)
            run.font.name = T.font

        # 카테고리 축
        cax = chart.category_axis
        cax.tick_labels.font.size = Pt(10)
        cax.tick_labels.font.color.rgb = T.rgb(T.black)
        cax.tick_labels.font.name = T.font
        cax.format.line.color.rgb = T.rgb(T.light_gray)
        cax.format.line.width = Pt(0.5)

        _hide_value_axis(chart)
        self._pgnum(slide, pg, tot)

    # ── 파이 / 도넛 ──
    def add_pie(self, labels, pcts, freqs, n, title, pg, tot, donut=False):
        slide = self._slide()
        self._title(slide, title, f"(n={n:,})")
        self._unit(slide, "%, 명")

        cd = CategoryChartData()
        cd.categories = labels
        cd.add_series(' ', pcts)

        ct = XL_CHART_TYPE.DOUGHNUT if donut else XL_CHART_TYPE.PIE
        cf = slide.shapes.add_chart(
            ct, Inches(1.5), Inches(1.2), Inches(10), Inches(5.5), cd
        )
        chart = cf.chart
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(11)
        chart.legend.font.color.rgb = T.rgb(T.black)
        chart.legend.font.name = T.font

        # 기본 라벨 설정
        plot = chart.plots[0]
        plot.has_data_labels = True
        dl = plot.data_labels
        dl.font.size = Pt(11)
        dl.font.color.rgb = T.rgb(T.white)
        dl.font.bold = True
        dl.font.name = T.font
        dl.number_format = '0.0"%"'

        # 색상
        colors = T.colors
        series = plot.series[0]
        for i, point in enumerate(series.points):
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = T.rgb(colors[i % len(colors)])

        # 커스텀 라벨
        for i in range(len(labels)):
            pt_dl = series.points[i].data_label
            tf = pt_dl.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = f"{pcts[i]:.1f}% ({int(freqs[i])})"
            run.font.size = Pt(11)
            run.font.color.rgb = T.rgb(T.white)
            run.font.bold = True
            run.font.name = T.font

        self._pgnum(slide, pg, tot)

    # ── 레이더 ──
    def add_radar(self, labels, values, title, pg, tot):
        slide = self._slide()
        self._title(slide, title, "(5점 만점 기준)")
        self._unit(slide, "점")

        cd = CategoryChartData()
        cd.categories = labels
        cd.add_series(' ', values)

        cf = slide.shapes.add_chart(
            XL_CHART_TYPE.RADAR_MARKERS,
            Inches(2.0), Inches(1.2), Inches(9), Inches(5.5), cd
        )
        chart = cf.chart
        chart.has_legend = False

        ser = chart.series[0]
        ser.format.line.color.rgb = T.rgb(T.accent)
        ser.format.line.width = Pt(2.5)
        ser.marker.style = 8
        ser.marker.size = 10
        ser.marker.format.fill.solid()
        ser.marker.format.fill.fore_color.rgb = T.rgb(T.accent3)

        plot = chart.plots[0]
        plot.has_data_labels = True
        dl = plot.data_labels
        dl.font.size = Pt(10)
        dl.font.color.rgb = T.rgb(T.accent)
        dl.font.bold = True
        dl.font.name = T.font
        dl.number_format = '0.00'

        cax = chart.category_axis
        cax.tick_labels.font.size = Pt(11)
        cax.tick_labels.font.color.rgb = T.rgb(T.black)
        cax.tick_labels.font.name = T.font

        va = chart.value_axis
        va.tick_labels.font.size = Pt(2)
        va.tick_labels.font.color.rgb = T.rgb(T.white)

        self._pgnum(slide, pg, tot)

    # ── 누적 막대 ──
    def add_stacked(self, items, data, univs, title, pg, tot):
        slide = self._slide()
        self._title(slide, title, f"(상위 {len(items)}개)")
        self._unit(slide, "명")

        r_items = list(reversed(items))
        cd = CategoryChartData()
        cd.categories = r_items
        for u in univs:
            cd.add_series(u, list(reversed(data[u])))

        cf = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_STACKED,
            Inches(0.5), Inches(1.0), Inches(12.3), Inches(5.9), cd
        )
        chart = cf.chart
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.TOP
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(10)
        chart.legend.font.color.rgb = T.rgb(T.black)
        chart.legend.font.name = T.font

        colors = T.colors
        for si, s in enumerate(chart.series):
            s.format.fill.solid()
            s.format.fill.fore_color.rgb = T.rgb(colors[si % len(colors)])

        plot = chart.plots[0]
        plot.has_data_labels = True
        dl = plot.data_labels
        dl.font.size = Pt(8)
        dl.font.color.rgb = T.rgb(T.white)
        dl.font.name = T.font
        dl.number_format = '#,##0'
        dl.label_position = XL_LABEL_POSITION.CENTER

        cax = chart.category_axis
        cax.tick_labels.font.size = Pt(9)
        cax.tick_labels.font.color.rgb = T.rgb(T.black)
        cax.tick_labels.font.name = T.font
        cax.format.line.color.rgb = T.rgb(T.light_gray)
        cax.format.line.width = Pt(0.5)

        _hide_value_axis(chart)
        self._pgnum(slide, pg, tot)

    def save(self, path):
        os.makedirs(os.path.dirname(path), exist_ok=True)
        self.prs.save(path)


# ═══════════════════════════════════════════════════════
# 메인 + 검증
# ═══════════════════════════════════════════════════════

def build(xlsx_path, out_path):
    print("=" * 60)
    print("  울산 미스매치 — PPTX 차트 생성 (v4 흰배경)")
    print("=" * 60)

    wb = openpyxl.load_workbook(xlsx_path, data_only=True)
    ws = wb['Sheet1']
    b = Builder()
    total = len(TABLES)
    ok = 0

    for c in TABLES:
        pg = ok + 1
        print(f"\n[표 {c.idx:2d}] {c.title}")
        try:
            if c.chart == "hbar":
                ls, ps, fs, n = get_bar_data(ws, c)
                flt = [(l, p, f) for l, p, f in zip(ls, ps, fs) if p > 0]
                if flt:
                    fl, fp, ff = zip(*flt)
                    b.add_hbar(list(fl), list(fp), list(ff), n, c.title, pg, total)
                    print(f"  → 수평막대: {len(fl)}항목"); ok += 1
                else: print("  → 데이터 없음")
            elif c.chart == "pie":
                if c.idx == 1:
                    ls, ps, vs, n = get_pie1(ws); fs = vs
                else:
                    ls, ps, fs, n = get_bar_data(ws, c)
                b.add_pie(ls, ps, fs, n, c.title, pg, total)
                print(f"  → 파이: {len(ls)}항목"); ok += 1
            elif c.chart == "donut":
                ls, ps, fs, n = get_bar_data(ws, c)
                b.add_pie(ls, ps, fs, n, c.title, pg, total, donut=True)
                print(f"  → 도넛: {len(ls)}항목"); ok += 1
            elif c.chart == "radar":
                ls, vs = get_radar(ws, c)
                if ls:
                    b.add_radar(ls, vs, c.title, pg, total)
                    print(f"  → 레이더: {len(ls)}항목"); ok += 1
            elif c.chart == "stacked":
                it, dt, un = get_stacked(ws, c)
                if it:
                    b.add_stacked(it, dt, un, c.title, pg, total)
                    print(f"  → 누적막대: {len(it)}항목"); ok += 1
        except Exception as e:
            print(f"  → 오류: {e}")
            import traceback; traceback.print_exc()

    wb.close()
    b.save(out_path)
    print(f"\n{'='*60}")
    print(f"  생성: {ok}장 → {out_path}")
    print(f"{'='*60}")

    # 검증
    print(f"\n[검증] 무결성 확인...")
    try:
        t = Presentation(out_path)
        n = len(t.slides)
        print(f"  ✓ 로드 성공, {n}장")
        for i, s in enumerate(t.slides):
            print(f"    슬라이드 {i+1}: {len(s.shapes)}개 shape")
        print(f"  ✓ 검증 통과")
    except Exception as e:
        print(f"  ✗ 실패: {e}")


if __name__ == "__main__":
    XLSX = r"d:\git_rk\project\25_121_ulsan\output\(OUTPUT)울산 일자리 미스매치_260205_1822.xlsx"
    ts = datetime.now().strftime('%y%m%d_%H%M')
    OUT  = rf"d:\git_rk\project\25_121_ulsan\output\wording\울산_미스매치_차트_PPT_{ts}.pptx"
    build(XLSX, OUT)
