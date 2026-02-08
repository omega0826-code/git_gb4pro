from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_indicator_slide():
    prs = Presentation()
    
    # 슬라이드 크기 설정 (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # 배경색 설정 (매우 연한 그레이 - 전문적인 느낌)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(252, 252, 252)
    
    # 제목 추가 (위치 및 폰트 개선)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "상권별 세부 지표 및 산출 로직 현황"
    p.font.bold = True
    p.font.size = Pt(24) # 크기 약간 축소
    p.font.name = 'Malgun Gothic'
    p.font.color.rgb = RGBColor(31, 73, 125) # 다크 블루
    
    # 데이터 정의
    headers = ["차원", "세부 지표", "산출 로직", "강남역", "선릉역", "마이스", "압구정", "청담", "가로수", "양재", "신논현", "역삼"]
    rows = [
        ["경쟁", "경쟁 점포수", "SUM(의원/피부과 점포)", "684", "108", "10", "242", "137", "142", "15", "17", "3"],
        ["고객", "타겟 매출(억)", "20/30/40대 매출 합계", "2,765", "573", "573", "739", "744", "836", "59.8", "21.4", "0.7"],
        ["", "여성 매출(억)", "전체 여성 매출 합계", "2,190", "457", "535", "514", "672", "619", "40.6", "11.5", "0.4"],
        ["인구", "유동인구(천)", "총 유동인구수(월)", "8,156", "2,143", "113", "5,371", "1,395", "2,008", "1,788", "1,308", "609"],
        ["", "타겟 유동(천)", "20-40대 유동인구 합", "5,869", "1,378", "79.8", "3,282", "851", "1,360", "961", "796", "449"],
        ["", "직장인구(명)", "총 직장인구수", "87,191", "11,262", "102,032", "59,959", "18,696", "15,957", "6,972", "1,661", "1,181"],
        ["", "상주인구(명)", "총 상주인구수", "6,248", "2,685", "34", "4,718", "2,738", "1,654", "5,763", "4,226", "742"],
        ["입지", "집객시설(개)", "상권 내 집객시설 수", "389", "44", "35", "152", "95", "77", "27", "8", "7"],
        ["", "평균 소득(만)", "월 평균 소득 금액", "445", "360", "701", "348", "408", "437", "374", "272", "411"],
        ["", "의료비 지출(백만)", "의료비 지출 총금액", "476", "284", "0", "455", "215", "192", "436", "331", "121"]
    ]

    # 표 추가 (크기 축소: 가로 1/2, 세로 2/3)
    num_rows = len(rows) + 1
    num_cols = len(headers)
    
    # 원본 가로 12.533 -> 6.266
    # 원본 세로 5.2 -> 3.466
    width = Inches(6.266)
    height = Inches(3.466)
    
    # 슬라이드 중앙 정렬
    left = (prs.slide_width - width) / 2
    top = Inches(1.8)
    
    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    table = table_shape.table
    
    # 컬럼 너비 비율 유지하며 축소
    total_ratio = 12.533 / 6.266
    orig_widths = [0.8, 1.2, 2.3, 0.9, 0.9, 1.0, 0.9, 0.9, 0.9, 0.9, 0.9, 0.9]
    for i, w in enumerate(orig_widths):
        table.columns[i].width = Inches(w / total_ratio)

    # 헤더 스타일링 (비즈니스 임팩트)
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(55, 96, 146) # 심플한 비즈니스 블루
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(8) # 크기에 맞춰 축소 (데이터가 많아 8pt가 적당)
        p.font.name = 'Malgun Gothic'
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
    # 데이터 행 스타일링
    for r_idx, row_data in enumerate(rows):
        for c_idx, val in enumerate(row_data):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(7) # 데이터 7pt로 조정하여 여백 확보
            p.font.name = 'Malgun Gothic'
            p.font.color.rgb = RGBColor(64, 64, 64)
            
            # 정렬
            if c_idx <= 2:
                p.alignment = PP_ALIGN.LEFT
            else:
                p.alignment = PP_ALIGN.CENTER
            
            # 교차 행 배경색 (매우 은은하게)
            if (r_idx + 1) % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(245, 248, 253)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
            
            # 차원 필드 텍스트 강조 (볼드)
            if c_idx == 0 and val:
                p.font.bold = True
                p.font.color.rgb = RGBColor(31, 73, 125)

    # 저장 (파일명 v2로 변경하여 덮어쓰기 방지)
    output_path = r"d:\git_gb4pro\data\gangnam\03_수요분석\상권별_지표_데이터_및_로직_v2.pptx"
    prs.save(output_path)
    print(f"PPT 고도화 완료 (v2): {output_path}")

if __name__ == "__main__":
    create_indicator_slide()
