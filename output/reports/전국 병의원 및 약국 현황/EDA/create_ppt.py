# -*- coding: utf-8 -*-
"""
서울시 피부과 EDA 보고서 - 편집 가능한 PPT 생성 (v2)
- 인사이트 워딩 추가
- 이미지 크기 PPT에 최적화
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path
from datetime import datetime

# 경로 설정
TIMESTAMP = datetime.now().strftime('%Y%m%d_%H%M')
BASE_DIR = Path(r'd:\git_gb4pro\output\reports\전국 병의원 및 약국 현황\EDA')
IMAGE_DIR = BASE_DIR / 'EDA_Dermatology_20260205_0424'
OUTPUT_DIR = BASE_DIR / 'PPT'
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
OUTPUT_FILE = OUTPUT_DIR / f'서울시_피부과_EDA_{TIMESTAMP}.pptx'

# 프레젠테이션 생성 (와이드스크린 16:9)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_title_slide(title, subtitle):
    """표지 슬라이드"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 배경
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(44, 62, 80)
    shape.line.fill.background()
    
    # 제목
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # 부제목
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(12.3), Inches(1))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(24)
    p2.font.color.rgb = RGBColor(189, 195, 199)
    p2.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(title, image_path, insight):
    """분석 슬라이드 - 이미지 + 인사이트"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 제목
    txBox = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(44, 62, 80)
    
    # 구분선
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(0.85), Inches(12.5), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(52, 152, 219)
    line.line.fill.background()
    
    # 이미지 (슬라이드에 맞게 크기 조정)
    if Path(image_path).exists():
        # 이미지 영역: 왼쪽 70%
        slide.shapes.add_picture(str(image_path), Inches(0.3), Inches(1.0), width=Inches(8.8))
    
    # 인사이트 박스 (오른쪽 30%)
    insight_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                          Inches(9.3), Inches(1.0), 
                                          Inches(3.7), Inches(6.2))
    insight_box.fill.solid()
    insight_box.fill.fore_color.rgb = RGBColor(248, 249, 250)
    insight_box.line.color.rgb = RGBColor(52, 152, 219)
    insight_box.line.width = Pt(2)
    
    # 인사이트 제목
    txBox_insight_title = slide.shapes.add_textbox(Inches(9.5), Inches(1.2), Inches(3.3), Inches(0.4))
    tf_title = txBox_insight_title.text_frame
    p_title = tf_title.paragraphs[0]
    p_title.text = "💡 Key Insight"
    p_title.font.size = Pt(16)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(52, 152, 219)
    
    # 인사이트 내용
    txBox_insight = slide.shapes.add_textbox(Inches(9.5), Inches(1.7), Inches(3.3), Inches(5.3))
    tf_insight = txBox_insight.text_frame
    tf_insight.word_wrap = True
    
    for line_text in insight.split('\n'):
        p_insight = tf_insight.add_paragraph()
        p_insight.text = line_text
        p_insight.font.size = Pt(13)
        p_insight.font.color.rgb = RGBColor(52, 73, 94)
        p_insight.space_after = Pt(8)
    
    return slide

def add_summary_slide(title, stats):
    """요약 슬라이드"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 제목
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(44, 62, 80)
    
    # 통계 카드
    card_width = Inches(2.8)
    card_height = Inches(2.2)
    start_x = Inches(0.5)
    start_y = Inches(1.8)
    gap = Inches(0.3)
    
    colors = [
        RGBColor(231, 76, 60),
        RGBColor(52, 152, 219),
        RGBColor(46, 204, 113),
        RGBColor(155, 89, 182)
    ]
    
    for i, (label, value, desc) in enumerate(stats):
        x = start_x + (card_width + gap) * i
        
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, start_y, card_width, card_height)
        card.fill.solid()
        card.fill.fore_color.rgb = colors[i % len(colors)]
        card.line.fill.background()
        
        # 값
        txBox = slide.shapes.add_textbox(x, start_y + Inches(0.3), card_width, Inches(0.8))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # 라벨
        txBox2 = slide.shapes.add_textbox(x, start_y + Inches(1.2), card_width, Inches(0.4))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = label
        p2.font.size = Pt(16)
        p2.font.bold = True
        p2.font.color.rgb = RGBColor(255, 255, 255)
        p2.alignment = PP_ALIGN.CENTER
        
        # 설명
        txBox3 = slide.shapes.add_textbox(x, start_y + Inches(1.6), card_width, Inches(0.5))
        tf3 = txBox3.text_frame
        p3 = tf3.paragraphs[0]
        p3.text = desc
        p3.font.size = Pt(11)
        p3.font.color.rgb = RGBColor(255, 255, 255)
        p3.alignment = PP_ALIGN.CENTER
    
    # 핵심 요약 박스
    summary_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                          Inches(0.5), Inches(4.5), 
                                          Inches(12.3), Inches(2.5))
    summary_box.fill.solid()
    summary_box.fill.fore_color.rgb = RGBColor(44, 62, 80)
    summary_box.line.fill.background()
    
    txBox_sum = slide.shapes.add_textbox(Inches(0.8), Inches(4.8), Inches(11.7), Inches(2))
    tf_sum = txBox_sum.text_frame
    tf_sum.word_wrap = True
    
    summary_points = [
        "✓ 서울시 피부과 의료기관의 98.5%가 의원급으로 소규모 전문 클리닉 중심",
        "✓ 강남구에 전체의 23.7% 집중 → 고소득층 밀집 지역 선호 경향",
        "✓ 최근 5년 신규 개원 비율 상승 → 피부과 시장 지속 성장 확인"
    ]
    
    for pt in summary_points:
        p_sum = tf_sum.add_paragraph()
        p_sum.text = pt
        p_sum.font.size = Pt(18)
        p_sum.font.color.rgb = RGBColor(255, 255, 255)
        p_sum.space_after = Pt(12)
    
    return slide

# 슬라이드 생성
print("PPT 슬라이드 생성 중...")

# 1. 표지
add_title_slide(
    "서울시 피부과 EDA 분석",
    "의원/병원/종합병원/상급종합 대상 | 데이터 기준: 2025년 12월"
)

# 2. 요약 슬라이드
add_summary_slide("📊 분석 핵심 지표", [
    ("총 기관 수", "4,853개", "서울시 전체 피부과"),
    ("강남구 집중도", "23.7%", "1,150개 / 전국 1위"),
    ("의원급 비중", "98.5%", "4,781개 / 소규모 중심"),
    ("평균 전문의", "1.30명", "기관당 의과전문의")
])

# 3-11. 분석 슬라이드 + 인사이트
slides_config = [
    ("기관 유형별 분포", "01_type_distribution.png",
     "• 의원급이 98.5%로 압도적\n• 병원급 이상은 1.5%\n• 소규모 전문 클리닉이 피부과 시장의 주류\n• 대형 병원보다 접근성 높은 의원 선호"),
    
    ("자치구별 분포 (Top 10)", "02_district_distribution.png",
     "• 강남구 1,150개로 압도적 1위\n• 서초구(451), 송파구(287) 순\n• 강남 3구에 38.9% 집중\n• 고소득층 밀집 지역과 상관관계"),
    
    ("설립구분별 분포", "03_establish_distribution.png",
     "• 개인 설립이 대다수\n• 법인/재단 설립 비율 미미\n• 개인 사업자 중심의 창업 시장\n• 진입장벽이 비교적 낮음"),
    
    ("연도별 개설 추이", "04_year_distribution.png",
     "• 2000년대 이후 급성장\n• 2010년대 개원 러시\n• 최근 5년간 신규 개원 증가\n• 피부과 시장 지속 성장 중"),
    
    ("종별 전문의 분포", "05_specialist_analysis.png",
     "• 의원급 평균 1.3명 전문의\n• 종합병원 이상은 다수 전문의\n• 1인 전문의 클리닉이 주류\n• 전문의 확보가 경쟁력"),
    
    ("병상 규모 현황", "06_bed_analysis.png",
     "• 대부분 외래 중심 운영\n• 입원 병상 보유 기관 소수\n• 피부과는 외래 진료 특화\n• 수술/입원 필요 시 상위 병원 의뢰"),
    
    ("병원당 평균 병상수", "06b_bed_per_hospital.png",
     "• 종합병원급만 유의미한 병상\n• 의원급은 대부분 병상 없음\n• 피부과 특성상 입원 수요 낮음\n• 외래 회전율이 수익의 핵심"),
    
    ("주요 병행 진료과목", "07_department_analysis.png",
     "• 성형외과 병행 최다\n• 내과, 가정의학과 순\n• 피부미용 + 성형 시너지\n• 복합 진료과목 운영 트렌드"),
    
    ("병원 연령대 분포", "08_age_size_analysis.png",
     "• 중견(10-20년) 기관 최다\n• 신규(5년 미만) 비율 상승\n• 노포(30년+) 기관 희소\n• 세대교체 진행 중"),
]

for title, img_file, insight in slides_config:
    img_path = IMAGE_DIR / img_file
    add_content_slide(title, img_path, insight)

# 저장
prs.save(OUTPUT_FILE)
print(f"PPT 저장 완료: {OUTPUT_FILE}")
